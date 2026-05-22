"""
新媒体PPT生成工具脚本
用法:
    python make_ppt.py --brand trumpchi --type weekly --data data.json --output report.pptx
    python make_ppt.py --brand audi --type monthly --output monthly.pptx
"""

import argparse
import json
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# 品牌配色方案
# ============================================================
BRAND_COLORS = {
    "trumpchi": {
        "primary": RGBColor(0, 82, 204),
        "secondary": RGBColor(230, 242, 255),
        "accent": RGBColor(0, 61, 153),
        "name": "广汽传祺",
    },
    "audi": {
        "primary": RGBColor(0, 0, 0),
        "secondary": RGBColor(192, 192, 192),
        "accent": RGBColor(255, 215, 0),
        "name": "上汽奥迪",
    },
    "hyper": {
        "primary": RGBColor(106, 13, 173),
        "secondary": RGBColor(240, 230, 255),
        "accent": RGBColor(255, 215, 0),
        "name": "广汽昊铂",
    },
    "aion": {
        "primary": RGBColor(0, 168, 107),
        "secondary": RGBColor(232, 245, 233),
        "accent": RGBColor(0, 122, 94),
        "name": "广汽埃安",
    },
}


# ============================================================
# 页面构建函数
# ============================================================
def _add_title_bar(slide, colors):
    """添加品牌色顶部装饰条"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors["primary"]
    bar.line.fill.background()


def _add_slide_title(slide, title, colors):
    """添加幻灯片标题"""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = colors["primary"]


def create_cover(prs, brand, data):
    """创建封面页"""
    colors = BRAND_COLORS[brand]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_bar(slide, colors)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.2))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = data.get("title", f"{colors['name']} 新媒体运营报告")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors["primary"]

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(8.4), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.add_paragraph()
    p.text = data.get("subtitle", "")
    p.font.size = Pt(24)
    p.font.color.rgb = colors["accent"]

    footer = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.5))
    tf = footer.text_frame
    p = tf.add_paragraph()
    date_str = data.get("date", datetime.now().strftime("%Y-%m-%d"))
    p.text = f"报告日期: {date_str}"
    if data.get("presenter"):
        p.text += f"  |  汇报人: {data['presenter']}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)


def create_metrics_slide(prs, brand, data):
    """创建核心数据指标页"""
    colors = BRAND_COLORS[brand]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_bar(slide, colors)
    _add_slide_title(slide, data.get("section_title", "核心数据概览"), colors)

    metrics = data.get("metrics", [])
    card_w = Inches(2.1)
    gap = Inches(0.2)
    start_x = Inches(0.5)
    y = Inches(1.5)

    for idx, m in enumerate(metrics[:4]):
        x = start_x + idx * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(1.3)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors["secondary"]
        card.line.fill.background()

        lb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05), card_w - Inches(0.2), Inches(0.3))
        p = lb.text_frame.add_paragraph()
        p.text = m.get("label", "")
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(128, 128, 128)

        vb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.35), card_w - Inches(0.2), Inches(0.5))
        p = vb.text_frame.add_paragraph()
        p.text = str(m.get("value", ""))
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]

        cb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.9), card_w - Inches(0.2), Inches(0.3))
        p = cb.text_frame.add_paragraph()
        change = m.get("change", "")
        p.text = change
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0, 128, 0) if change.startswith("+") else (
            RGBColor(192, 0, 0) if change.startswith("-") else RGBColor(128, 128, 128)
        )

    # 核心洞察
    insights_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2), Inches(9), Inches(3.5)
    )
    tf = insights_box.text_frame
    tf.word_wrap = True

    if data.get("insights"):
        p = tf.add_paragraph()
        p.text = "核心洞察"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.space_after = Pt(10)

        for insight in data["insights"]:
            p = tf.add_paragraph()
            p.text = f"  {chr(9654)} {insight}"
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(6)


def create_table_slide(prs, brand, data):
    """创建数据表格页"""
    colors = BRAND_COLORS[brand]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_bar(slide, colors)
    _add_slide_title(slide, data.get("section_title", "数据明细"), colors)

    headers = data.get("headers", [])
    rows = data.get("rows", [])

    if not headers or not rows:
        return

    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), Inches(1.3),
        Inches(9), Inches(0.4) * min(num_rows, 13)
    ).table

    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = colors["primary"]
        for p in c.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.text = str(val)
            if ri % 2 == 1:
                c.fill.solid()
                c.fill.fore_color.rgb = colors["secondary"]
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER


def create_summary(prs, brand, data):
    """创建总结页"""
    colors = BRAND_COLORS[brand]
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_bar(slide, colors)
    _add_slide_title(slide, "总结与下一步", colors)

    # 核心观点
    box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "核心结论"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = colors["primary"]

    for pt in data.get("key_points", []):
        p = tf.add_paragraph()
        p.text = f"  {chr(9679)} {pt}"
        p.font.size = Pt(14)
        p.space_after = Pt(8)

    # 下一步
    box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "下一步行动"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = colors["primary"]

    for idx, step in enumerate(data.get("next_steps", []), 1):
        p = tf.add_paragraph()
        p.text = f"  {idx}. {step}"
        p.font.size = Pt(14)
        p.space_after = Pt(8)


# ============================================================
# 主流程
# ============================================================
TYPE_HANDLERS = {
    "weekly": [create_cover, create_metrics_slide, create_table_slide, create_summary],
    "monthly": [create_cover, create_metrics_slide, create_table_slide, create_table_slide, create_summary],
    "analysis": [create_cover, create_table_slide, create_metrics_slide, create_summary],
    "strategy": [create_cover, create_table_slide, create_summary],
    "activity": [create_cover, create_metrics_slide, create_table_slide, create_summary],
}


def generate(brand, ppt_type, data, output_path):
    """主生成函数"""
    if brand not in BRAND_COLORS:
        raise ValueError(f"不支持的品牌: {brand}，可选: {list(BRAND_COLORS.keys())}")

    if ppt_type not in TYPE_HANDLERS:
        raise ValueError(f"不支持的类型: {ppt_type}，可选: {list(TYPE_HANDLERS.keys())}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 按 type 决定包含哪些 slide
    sections = data.get("sections", [])
    handlers = TYPE_HANDLERS[ppt_type]

    for i, handler in enumerate(handlers):
        section_data = sections[i] if i < len(sections) else {}
        # 合并顶层 data 和 section 数据
        merged = {**data, **section_data}
        handler(prs, brand, merged)

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="新媒体PPT生成工具")
    parser.add_argument("--brand", required=True,
                        choices=["trumpchi", "audi", "hyper", "aion"],
                        help="品牌标识")
    parser.add_argument("--type", dest="ppt_type", default="weekly",
                        choices=["weekly", "monthly", "analysis", "strategy", "activity"],
                        help="PPT类型")
    parser.add_argument("--data", help="JSON数据文件路径")
    parser.add_argument("--output", required=True, help="输出PPT文件路径")
    args = parser.parse_args()

    # 加载数据
    if args.data:
        with open(args.data, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = {
            "title": f"{BRAND_COLORS[args.brand]['name']} 新媒体运营报告",
            "subtitle": f"{args.ppt_type} 报告",
            "date": datetime.now().strftime("%Y-%m-%d"),
            "sections": [],
        }

    output = generate(args.brand, args.ppt_type, data, args.output)
    print(f"PPT 已生成: {output}")


if __name__ == "__main__":
    main()
